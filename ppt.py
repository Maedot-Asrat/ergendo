from pptx import Presentation

prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Program Blueprint Briefing"
subtitle.text = "Trive Club - Preparation & Planning"

# Phases slides
phases = [
    ("Phase One", "Finish designing the training program and set dates.\nFinish setting capstone projects.\n\nGoal: Detailed roadmap combining training and capstones.\nDate: October 14"),
    ("Phase Two", "Set up evaluation criteria for mentors, mentees, trainees, apprentices (attendance, participation).\n\nGoal: Clear methods to gauge everyone's performance.\nDate: October 14"),
    ("Phase Three", "Assign mentors to specific training sessions.\nOption 1: Let them choose.\nOption 2: Assign sessions.\n\nGoal: Ensure all sessions are booked.\nDate: October 16"),
    ("Phase Four", "Mentor orientation: expectations, rules, reveal program, apprentices, support.\n\nGoal: Mentors understand goals and role.\nDate: October 18"),
    ("Phase Five", "Choose X students from 49 applicants.\nConfirm interest.\nPair with mentors.\nReplace uninterested students.\n\nGoal: Right people get mentee spots.\nDate: Oct 19-22"),
    ("Phase Six", "Interview and accept apprentices.\nAssign them as assistants based on strengths.\nConnect apprentices with mentors.\n\nGoal: Strong support for each session.\nDate: Oct 19-22"),
    ("Phase Seven", "Opening Ceremony\n\nDate: October 25")
]

for title_text, content in phases:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = content

file_path = "/mnt/data/trive_club_briefing.pptx"
prs.save(file_path)
file_path
